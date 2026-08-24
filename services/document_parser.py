from io import BytesIO
from pathlib import Path

import fitz
from docx import Document
from pptx import Presentation

from config import MAX_UPLOAD_BYTES, SUPPORTED_EXTENSIONS


class DocumentError(ValueError):
    pass


def parse_document(filename: str, content: bytes) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        raise DocumentError(f"지원하지 않는 형식입니다: {suffix or '확장자 없음'}")
    if not content:
        raise DocumentError("빈 파일은 업로드할 수 없습니다.")
    if len(content) > MAX_UPLOAD_BYTES:
        raise DocumentError(f"파일은 {MAX_UPLOAD_BYTES // 1024 // 1024}MB 이하여야 합니다.")

    try:
        if suffix == ".pdf":
            text = _parse_pdf(content)
        elif suffix == ".docx":
            text = _parse_docx(content)
        elif suffix == ".pptx":
            text = _parse_pptx(content)
        else:
            text = _decode_text(content)
    except DocumentError:
        raise
    except Exception as exc:
        raise DocumentError(f"{filename}을 읽지 못했습니다: {exc}") from exc

    text = "\n".join(line.rstrip() for line in text.splitlines()).strip()
    if not text:
        raise DocumentError("추출할 텍스트가 없습니다. 스캔 PDF는 OCR 처리 후 다시 업로드해 주세요.")
    return text


def _parse_pdf(content: bytes) -> str:
    with fitz.open(stream=content, filetype="pdf") as document:
        return "\n\n".join(page.get_text("text") for page in document)


def _parse_docx(content: bytes) -> str:
    document = Document(BytesIO(content))
    paragraphs = [p.text for p in document.paragraphs]
    table_cells = [cell.text for table in document.tables for row in table.rows for cell in row.cells]
    return "\n".join(paragraphs + table_cells)


def _parse_pptx(content: bytes) -> str:
    presentation = Presentation(BytesIO(content))
    lines: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        if slide_text:
            lines.append(f"[슬라이드 {slide_number}]\n" + "\n".join(slide_text))
    return "\n\n".join(lines)


def _decode_text(content: bytes) -> str:
    for encoding in ("utf-8-sig", "cp949"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            pass
    raise DocumentError("텍스트 인코딩은 UTF-8 또는 CP949여야 합니다.")

